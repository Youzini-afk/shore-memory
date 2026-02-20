import concurrent.futures
import json
import mimetypes
import os
import shutil
import subprocess
import threading
from datetime import datetime
from typing import List, Optional

# 文档处理导入
try:
    import docx
except ImportError:
    docx = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import pypdf
except ImportError:
    pypdf = None


def find_es_executable() -> Optional[str]:
    """在常用路径中查找 'es.exe' 可执行文件。"""
    # 1. 检查 PATH 环境变量
    path_in_env = shutil.which("es")
    if path_in_env:
        return path_in_env

    # 2. 检查常用安装目录和本地工具
    current_dir = os.path.dirname(os.path.abspath(__file__))
    # 假设 file_search.py 位于 backend/tools/，我们在 backend/tools/ 和 backend/bin/ 中查找
    # 同时检查标准程序文件
    possible_paths = [
        os.path.join(current_dir, "es.exe"),
        os.path.join(
            os.path.dirname(current_dir), "bin", "es.exe"
        ),  # backend/bin/es.exe
        "C:\\Program Files\\Everything\\es.exe",
        "C:\\Program Files (x86)\\Everything\\es.exe",
    ]

    for path in possible_paths:
        if os.path.exists(path):
            return path

    return None


def fast_search_fallback(
    query: str, root_dir: str, limit: int = 50, timeout: float = 10.0
) -> List[str]:
    """
    带超时的多线程递归文件搜索。
    并行扫描 root_dir 的顶级目录。
    """
    results = []
    dirs_to_scan = []
    stop_event = threading.Event()

    try:
        # 获取顶级条目
        with os.scandir(root_dir) as it:
            for entry in it:
                if entry.name.startswith(".") or entry.name.startswith("$"):
                    continue

                if entry.is_dir():
                    dirs_to_scan.append(entry.path)
                elif entry.is_file() and query.lower() in entry.name.lower():
                    results.append(entry.path)
    except (PermissionError, OSError):
        pass

    if len(results) >= limit:
        return results

    # 线程辅助函数
    def scan_tree(path):
        local_res = []
        try:
            for root, dirs, files in os.walk(path):
                # 检查是否应该停止
                if stop_event.is_set():
                    return local_res

                # 优化：原地修改 dirs 以跳过隐藏目录
                dirs[:] = [
                    d for d in dirs if not d.startswith(".") and not d.startswith("$")
                ]

                for file in files:
                    if query.lower() in file.lower():
                        local_res.append(os.path.join(root, file))
                        if len(local_res) >= limit:
                            return local_res
        except Exception:
            pass
        return local_res

    # 使用 ThreadPoolExecutor 进行并行扫描
    with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
        future_to_path = {executor.submit(scan_tree, d): d for d in dirs_to_scan}

        try:
            # 等待结果（全局超时）
            for future in concurrent.futures.as_completed(
                future_to_path, timeout=timeout
            ):
                try:
                    data = future.result()
                    results.extend(data)
                    if len(results) >= limit:
                        stop_event.set()  # 通知其他线程停止
                        break
                except Exception:
                    pass
        except concurrent.futures.TimeoutError:
            print(f"[FileSearch] Parallel search timed out after {timeout}s")
            stop_event.set()  # 停止任何正在进行的工作

    return results[:limit]


def search_files(query: str, limit: int = 50) -> str:
    """
    在计算机上搜索文件。
    如果可用，首先尝试使用 'Everything' 命令行工具 (es.exe) 进行即时全局搜索。
    如果不可用，则回退到使用并行搜索（带超时）在用户主目录中搜索。

    Args:
        query: 要搜索的文件名或模式。
        limit: 返回的最大结果数。

    Returns:
        包含文件路径列表的 JSON 字符串。
    """
    results = []

    es_path = find_es_executable()

    # 1. 尝试 'Everything' (es.exe)
    if es_path:
        try:
            # -n <limit> 限制结果
            # -utf8 强制输出 UTF-8
            cmd = [es_path, query, "-n", str(limit), "-utf8"]

            # 使用 STARTUPINFO 在 Windows 上正确隐藏控制台窗口
            startupinfo = None
            if os.name == "nt":
                startupinfo = subprocess.STARTUPINFO()
                startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
                startupinfo.wShowWindow = subprocess.SW_HIDE  # 显式隐藏

            raw_output = subprocess.check_output(
                cmd,
                startupinfo=startupinfo,
                stderr=subprocess.STDOUT,  # 捕获 stderr 以便更好地调试
                stdin=subprocess.DEVNULL,
                creationflags=subprocess.CREATE_NO_WINDOW if os.name == "nt" else 0,
                timeout=10,  # [修复] 防止无限挂起
            )

            # 先尝试 UTF-8（因为传递了 -utf8），然后回退到系统编码
            try:
                output = raw_output.decode("utf-8")
            except UnicodeDecodeError:
                import locale

                system_enc = locale.getpreferredencoding()
                print(f"[FileSearch] UTF-8 decode failed, falling back to {system_enc}")
                output = raw_output.decode(system_enc, errors="ignore")

            lines = output.strip().split("\n")
            results = [line.strip() for line in lines if line.strip()]
            if results:
                final_results = results[:limit]
                return json.dumps(final_results, ensure_ascii=False)
        except (subprocess.CalledProcessError, FileNotFoundError, OSError) as e:
            print(f"[FileSearch] Error running es: {e}")
            if hasattr(e, "output") and e.output:
                print(
                    f"[FileSearch] Subprocess output: {e.output.decode('utf-8', errors='ignore')}"
                )
            pass  # 回退

    # 2. 回退：在用户主目录中使用并行搜索
    user_home = os.path.expanduser("~")
    # 使用 10s 超时进行回退，以保持 Pero 的响应性
    print(
        f"[FileSearch] 'es' not found. Falling back to parallel search in {user_home} (10s timeout)"
    )

    results = fast_search_fallback(query, user_home, limit, timeout=10.0)

    # 返回纯 JSON 供后端处理
    return json.dumps(results, ensure_ascii=False)


def read_file_content(file_path: str, max_length: int = 10000) -> str:
    """
    读取文件内容。支持文本文件、PDF、DOCX 和 PPTX。

    Args:
        file_path: 文件的绝对路径。
        max_length: 读取的最大字符数 (默认 10000)。

    Returns:
        文件内容字符串或错误消息。
    """
    try:
        # 检查文件是否存在
        if not os.path.exists(file_path):
            return f"错误: 未找到文件: {file_path}"

        # 检查是否为文件
        if not os.path.isfile(file_path):
            return f"错误: 路径不是文件: {file_path}"

        # 检查大小以避免读取超大文件
        file_size = os.path.getsize(file_path)
        if file_size > 10 * 1024 * 1024:  # 为了安全起见，硬限制 10MB
            return f"错误: 文件过大 ({file_size} bytes)。最大允许大小为 10MB。"

        ext = os.path.splitext(file_path)[1].lower()

        # 1. 处理 PDF
        if ext == ".pdf":
            if not pypdf:
                return "Error: pypdf is not installed. Cannot read PDF files."
            try:
                reader = pypdf.PdfReader(file_path)
                text = ""
                for page in reader.pages:
                    text += page.extract_text()
                    if len(text) > max_length:
                        break
                return text[:max_length] + (
                    "\n...[Content Truncated]..." if len(text) > max_length else ""
                )
            except Exception as e:
                return f"Error reading PDF: {str(e)}"

        # 2. 处理 DOCX
        elif ext == ".docx":
            if not docx:
                return "Error: python-docx is not installed. Cannot read DOCX files."
            try:
                doc = docx.Document(file_path)
                full_text = []
                for para in doc.paragraphs:
                    full_text.append(para.text)
                text = "\n".join(full_text)
                return text[:max_length] + (
                    "\n...[Content Truncated]..." if len(text) > max_length else ""
                )
            except Exception as e:
                return f"Error reading DOCX: {str(e)}"

        # 3. 处理 PPTX
        elif ext == ".pptx":
            if not Presentation:
                return "Error: python-pptx is not installed. Cannot read PPTX files."
            try:
                prs = Presentation(file_path)
                text_runs = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                text = "\n".join(text_runs)
                return text[:max_length] + (
                    "\n...[Content Truncated]..." if len(text) > max_length else ""
                )
            except Exception as e:
                return f"Error reading PPTX: {str(e)}"

        # 4. 处理纯文本（默认）
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read(max_length)
                if len(content) == max_length:
                    content += "\n...[Content Truncated]..."
                return content
        except UnicodeDecodeError:
            # 尝试几种常见编码
            for enc in ["gbk", "latin1"]:
                try:
                    with open(file_path, "r", encoding=enc) as f:
                        content = f.read(max_length)
                        if len(content) == max_length:
                            content += "\n...[Content Truncated]..."
                        return content
                except Exception:
                    continue
            return f"Error: Unable to decode file content. The file extension is {ext}, but it doesn't appear to be a standard text file or supported document format."

    except Exception as e:
        return f"Error reading file: {str(e)}"


def list_directory(path: str = None, dir_path: str = None) -> str:
    """
    列出给定目录中的文件和子目录。
    支持 'path' 和 'dir_path' 参数以提高灵活性。

    Args:
        path: 目录的绝对路径 (dir_path 的别名)。
        dir_path: 目录的绝对路径。

    Returns:
        目录内容的 JSON 字符串。
    """
    target_path = dir_path if dir_path else path
    if not target_path:
        target_path = "."  # 如果未提供，则默认为当前目录

    try:
        if not os.path.exists(target_path):
            return json.dumps({"error": f"Directory not found: {target_path}"})

        if not os.path.isdir(target_path):
            return json.dumps({"error": f"Path is not a directory: {target_path}"})

        items = []
        with os.scandir(target_path) as it:
            for entry in it:
                try:
                    item_type = "directory" if entry.is_dir() else "file"
                    items.append(
                        {"name": entry.name, "path": entry.path, "type": item_type}
                    )
                except OSError:
                    continue

        return json.dumps(items, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"Failed to list directory: {str(e)}"})


def get_file_info(file_path: str) -> str:
    """
    获取有关文件的元数据信息。

    Args:
        file_path: 文件的绝对路径。

    Returns:
        包含文件元数据的 JSON 字符串。
    """
    try:
        if not os.path.exists(file_path):
            return json.dumps({"error": f"未找到文件: {file_path}"})

        stat = os.stat(file_path)
        mime_type, _ = mimetypes.guess_type(file_path)

        info = {
            "name": os.path.basename(file_path),
            "path": file_path,
            "size": stat.st_size,
            "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "is_dir": os.path.isdir(file_path),
            "mime_type": mime_type or "unknown",
        }

        return json.dumps(info, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"Failed to get info: {str(e)}"})


def show_file_results(files: List[str]) -> str:
    """
    向用户显示包含找到文件列表的 UI。
    当你找到文件并希望将其展示给用户供选择时使用此功能。

    Args:
        files: 要显示的绝对文件路径列表。
    """
    # 我们返回指令给 LLM 以输出标签。
    # 前端仅解析助手消息中的标签。
    return f"UI Data prepared. To display the file list UI to the user, you MUST include this exact tag in your final response:\n<FILE_RESULTS>{json.dumps(files, ensure_ascii=False)}</FILE_RESULTS>"


# 工具定义
search_files_definition = {
    "type": "function",
    "function": {
        "name": "search_files",
        "description": "在整台计算机上按文件名搜索文件。当你需要根据文件名查找文件位置时使用此功能。(要在代码文件内部搜索，请改用 'code_search')。",
        "parameters": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "要搜索的文件名或关键字。",
                },
                "limit": {
                    "type": "integer",
                    "description": "返回的最大结果数 (默认 50)。",
                    "default": 50,
                },
            },
            "required": ["query"],
        },
    },
}

read_file_definition = {
    "type": "function",
    "function": {
        "name": "read_file_content",
        "description": "读取文档或文本文件的完整内容。最适合读取 .pdf, .docx, .pptx 或一般文本文件。(要跨多个代码文件搜索模式，请使用 'code_search')。",
        "parameters": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "要读取的文件的绝对路径。",
                },
                "max_length": {
                    "type": "integer",
                    "description": "读取的最大字符数 (默认 10000)。",
                    "default": 10000,
                },
            },
            "required": ["file_path"],
        },
    },
}

list_directory_definition = {
    "type": "function",
    "function": {
        "name": "list_directory",
        "description": "列出特定文件夹中的所有文件和子目录。用于探索文件结构。",
        "parameters": {
            "type": "object",
            "properties": {
                "dir_path": {
                    "type": "string",
                    "description": "要列出的目录的绝对路径。",
                }
            },
            "required": ["dir_path"],
        },
    },
}

get_file_info_definition = {
    "type": "function",
    "function": {
        "name": "get_file_info",
        "description": "获取关于文件或目录的元数据 (大小, 日期, 类型)。",
        "parameters": {
            "type": "object",
            "properties": {
                "file_path": {
                    "type": "string",
                    "description": "文件或目录的绝对路径。",
                }
            },
            "required": ["file_path"],
        },
    },
}

show_file_results_definition = {
    "type": "function",
    "function": {
        "name": "show_file_results",
        "description": "在 UI 覆盖层中向用户显示文件列表。在使用 search_files 找到文件后调用此功能。",
        "parameters": {
            "type": "object",
            "properties": {
                "files": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "要显示的文件路径列表。",
                }
            },
            "required": ["files"],
        },
    },
}
